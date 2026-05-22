import os
from io import BytesIO

import streamlit as st
from groq import Groq
from pptx import Presentation

# ----------------------------
# PAGE CONFIG
# ----------------------------
st.set_page_config(
    page_title="RAGHAV REALTY SEO Generator",
    page_icon="🏢",
    layout="wide"
)

# ----------------------------
# CUSTOM RAGHAV REALTY WEBSITE-INSPIRED STYLING
# ----------------------------
st.markdown("""
<style>
    html, body, [data-testid="stAppViewContainer"], .stApp {
        background: #f7f1e6 !important;
        background-image:
            radial-gradient(circle at top left, rgba(190, 143, 66, 0.18), transparent 30%),
            linear-gradient(180deg, #fff8ec 0%, #f4ead8 50%, #eadcc5 100%) !important;
        color: #2b2118 !important;
        font-family: Georgia, 'Times New Roman', serif !important;
    }

    [data-testid="stHeader"] {
        background: transparent !important;
    }

    [data-testid="stMainBlockContainer"], .block-container {
        background: rgba(255, 252, 245, 0.82) !important;
        border: 1px solid rgba(180, 132, 57, 0.24) !important;
        border-radius: 28px !important;
        padding: 2.2rem 2.4rem !important;
        max-width: 1120px !important;
        box-shadow: 0 24px 70px rgba(80, 55, 25, 0.12) !important;
    }

    .main-title {
        font-size: 2.9rem !important;
        font-weight: 800 !important;
        color: #8b5e2a !important;
        text-align: center !important;
        letter-spacing: 4px !important;
        text-transform: uppercase !important;
        margin-bottom: 0.25rem !important;
    }

    .subtitle {
        text-align: center !important;
        color: #6d5133 !important;
        margin-bottom: 2rem !important;
        letter-spacing: 1px !important;
        font-size: 1.08rem !important;
        font-family: Arial, sans-serif !important;
    }

    h1, h2, h3, h4, h5, h6, p, span, div, label, .stMarkdown {
        color: #2b2118 !important;
    }

    label {
        font-family: Arial, sans-serif !important;
        font-weight: 700 !important;
        color: #5b4228 !important;
    }

    .stTextInput input, .stTextArea textarea {
        background-color: #fffaf2 !important;
        color: #2b2118 !important;
        border: 1px solid rgba(139, 94, 42, 0.42) !important;
        border-radius: 14px !important;
        font-family: Arial, sans-serif !important;
    }

    [data-baseweb="select"] > div {
        background-color: #fffaf2 !important;
        color: #2b2118 !important;
        border: 1px solid rgba(139, 94, 42, 0.42) !important;
        border-radius: 14px !important;
        font-family: Arial, sans-serif !important;
    }

    [data-baseweb="select"] * {
        color: #2b2118 !important;
    }

    .stButton > button, .stDownloadButton > button {
        background: linear-gradient(135deg, #8b5e2a 0%, #c79b52 100%) !important;
        color: #fff8ec !important;
        border-radius: 999px !important;
        border: none !important;
        padding: 0.9rem 1.6rem !important;
        font-weight: 800 !important;
        width: 100% !important;
        font-family: Arial, sans-serif !important;
        letter-spacing: 0.5px !important;
        box-shadow: 0 12px 30px rgba(139, 94, 42, 0.24) !important;
    }

    .stTabs [data-baseweb="tab"] {
        background-color: #fff6e8 !important;
        color: #8b5e2a !important;
        border-radius: 999px !important;
        border: 1px solid rgba(139, 94, 42, 0.28) !important;
        padding: 12px 26px !important;
        font-weight: 800 !important;
        font-family: Arial, sans-serif !important;
    }

    .stTabs [aria-selected="true"] {
        background: #8b5e2a !important;
        color: #fff8ec !important;
    }
</style>
""", unsafe_allow_html=True)

# ----------------------------
# HEADER + LOGO
# ----------------------------
if os.path.exists("logo.png"):
    col_logo1, col_logo2, col_logo3 = st.columns([1, 1, 1])
    with col_logo2:
        st.image("logo.png", width=180)

st.markdown('<div class="main-title">RAGHAV REALTY</div>', unsafe_allow_html=True)
st.markdown('<div class="subtitle">AI-Powered SEO & Hashtag Generator for New Project Launches</div>', unsafe_allow_html=True)

# ----------------------------
# API CONFIG
# ----------------------------
client = Groq(api_key=st.secrets["GROQ_API_KEY"])

# ----------------------------
# PPT EXPORT FUNCTION
# ----------------------------
def create_ppt(project_name, seo_titles, meta_descriptions, keywords, hashtags, captions):
    prs = Presentation()

    slides = [
        ("SEO Titles", seo_titles),
        ("Meta Descriptions", meta_descriptions),
        ("SEO Keywords", keywords),
        ("Hashtags", hashtags),
        ("Captions & CTA", captions),
    ]

    for title, content in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"{project_name} - {title}"
        slide.placeholders[1].text = content

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# ----------------------------
# FORM
# ----------------------------
with st.form("seo_form"):
    col1, col2 = st.columns(2)

    with col1:
        project_name = st.text_input("Project Name")
        city = st.selectbox("City", ["Mumbai", "Pune", "Delhi", "Bangalore", "Jaipur", "Ahmedabad"])
        micro_market = st.text_input("Micro Market", placeholder="Example: Andheri East, Chembur, BKC")
        project_type = st.selectbox("Project Type", ["Residential", "Commercial", "Mixed"])

    with col2:
        configuration = st.multiselect(
            "Configuration",
            ["1 BHK", "2 BHK", "3 BHK", "4 BHK", "Office", "Retail", "Studio", "Penthouse"]
        )
        landmarks = st.text_input("Nearby Landmarks", placeholder="Example: Metro station, airport, school, mall")
        brand_positioning = st.selectbox("Brand Positioning", ["Luxury", "Premium", "Affordable", "Ultra Luxury"])

    usp = st.text_area("USP (Unique Selling Proposition)", placeholder="Example: Rooftop amenities, sea view, smart homes, excellent connectivity")

    submit = st.form_submit_button("Generate SEO & Hashtags")

# ----------------------------
# GENERATE OUTPUT
# ----------------------------
if submit:
    if not project_name or not micro_market or not usp:
        st.warning("Please fill Project Name, Micro Market, and USP for better output.")
    else:
        prompt = f"""
You are India's top luxury real estate SEO strategist and social media expert.

Your job is to generate HIGH-CONVERTING SEO content for premium real estate launches.

The content must:
- Improve Google organic rankings
- Improve Instagram reach and discoverability
- Improve LinkedIn engagement
- Target high-intent homebuyers and investors
- Use luxury and aspirational language
- Include strong local SEO keywords
- Focus on location-based search intent

PROJECT DETAILS:
Project Name: {project_name}
City: {city}
Micro Market: {micro_market}
Project Type: {project_type}
Configuration: {configuration}
Nearby Landmarks: {landmarks}
Brand Positioning: {brand_positioning}
USP: {usp}

Return the answer EXACTLY in this format:

SEO_TITLES:
Generate 5 premium SEO titles under 65 characters.

META_DESCRIPTIONS:
Generate 3 compelling SEO meta descriptions under 155 characters.

SEO_KEYWORDS:
Generate:
- 10 primary SEO keywords
- 10 long-tail keywords
- 5 buyer intent keywords

HASHTAGS:
Generate:
- 10 Instagram hashtags
- 10 luxury real estate hashtags
- 10 location-specific hashtags
- 5 investment-focused hashtags

CAPTIONS_AND_CTA:
Generate:
- 3 Instagram captions
- 3 Google Ads headlines
- 3 CTA lines

RULES:
- Use premium real estate vocabulary
- Avoid generic outputs
- Make hashtags highly searchable
- Use local SEO naturally
- Prioritize discoverability
- Keep output clean and easy to copy
"""

        with st.spinner("Generating premium SEO and hashtags..."):
            try:
                response = client.chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.7,
                    max_tokens=1800
                )

                result = response.choices[0].message.content

                def extract_section(text, start, end=None):
                    if start not in text:
                        return ""
                    section = text.split(start, 1)[1]
                    if end and end in section:
                        section = section.split(end, 1)[0]
                    return section.strip()

                seo_titles = extract_section(result, "SEO_TITLES:", "META_DESCRIPTIONS:")
                meta_descriptions = extract_section(result, "META_DESCRIPTIONS:", "SEO_KEYWORDS:")
                keywords = extract_section(result, "SEO_KEYWORDS:", "HASHTAGS:")
                hashtags = extract_section(result, "HASHTAGS:", "CAPTIONS_AND_CTA:")
                captions = extract_section(result, "CAPTIONS_AND_CTA:")

                st.success("SEO content generated successfully!")

                tab1, tab2, tab3, tab4, tab5 = st.tabs([
                    "SEO Titles",
                    "Meta Descriptions",
                    "Keywords",
                    "Hashtags",
                    "Captions & CTA"
                ])

                with tab1:
                    st.text_area("Copy SEO Titles", seo_titles, height=220)

                with tab2:
                    st.text_area("Copy Meta Descriptions", meta_descriptions, height=220)

                with tab3:
                    st.text_area("Copy SEO Keywords", keywords, height=320)

                with tab4:
                    st.text_area("Copy Hashtags", hashtags, height=320)

                with tab5:
                    st.text_area("Copy Captions & CTA", captions, height=320)

                ppt_file = create_ppt(
                    project_name,
                    seo_titles,
                    meta_descriptions,
                    keywords,
                    hashtags,
                    captions
                )

                st.download_button(
                    label="Download as PPT",
                    data=ppt_file,
                    file_name=f"{project_name.replace(' ', '_')}_SEO_Hashtags.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

            except Exception as e:
                st.error("Groq API Error:")
                st.code(str(e))


